import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement

# ---------------------------------------------------------
# CONSTANTS & CONFIGURATION
# ---------------------------------------------------------
OUTPUT_PATH = "SyncTrip_new.pptx"
ERD_IMAGE_PATH = "synctrip_clean_erd_diagram.jpg"

# Color Palette (White & Green Theme)
COLOR_WHITE = RGBColor(255, 255, 255)         # #FFFFFF
COLOR_LIGHT_GREEN = RGBColor(211, 245, 221)    # #D3F5DD (Soft green fill)
COLOR_GREEN = RGBColor(0, 188, 112)           # #00BC70 (Theme Primary Green)
COLOR_DARK_TEXT = RGBColor(26, 26, 26)         # #1A1A1A (Primary text)
COLOR_MUTED_TEXT = RGBColor(110, 110, 110)     # #6E6E6E (Subtext)
COLOR_YELLOW = RGBColor(255, 242, 117)         # #FFF275 (Yellow highlight chip)
COLOR_BORDER_GRAY = RGBColor(210, 225, 215)    # Grayish-green soft border

# Fonts (Latin & Korean pairs to ensure perfect rendering in all environments)
FONT_TITLE_EN = "KIMMBold"
FONT_TITLE_KO = "한국기계연구원_Bold"

FONT_BODY_EN = "A2Z 5 Medium"
FONT_BODY_KO = "에이투지체 5 Medium"

FONT_BODY_BOLD_EN = "A2Z 7 Bold"
FONT_BODY_KO_BOLD = "에이투지체 7 Bold"

FONT_BODY_REG_EN = "A2Z 4 Regular"
FONT_BODY_KO_REG = "에이투지체 4 Regular"

FONT_BODY_LIGHT_EN = "A2Z 3 Light"
FONT_BODY_KO_LIGHT = "에이투지체 3 Light"

# ---------------------------------------------------------
# HELPER FUNCTIONS
# ---------------------------------------------------------
def apply_font_to_run(run, font_en, font_ko, size_pt, color_rgb, bold=False):
    """Applies font formatting to a text run, forcing Latin and East Asian font matching."""
    run.font.name = font_en
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color_rgb
    run.font.bold = bold
    
    rPr = run.font._element
    for child in list(rPr):
        if child.tag.endswith('ea'):
            rPr.remove(child)
    
    ea = OxmlElement('a:ea')
    ea.set('typeface', font_ko)
    rPr.append(ea)

def add_header(slide, sub_title_text, title_text, page_num):
    """Creates a consistent premium header template on the slide."""
    # Top thin line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_GREEN
    line.line.fill.background()
    
    # Sub title
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(5.0), Inches(0.4))
    tf_sub = sub_box.text_frame
    tf_sub.margin_left = Inches(0)
    tf_sub.margin_top = Inches(0)
    p_sub = tf_sub.paragraphs[0]
    run_sub = p_sub.add_run()
    run_sub.text = sub_title_text
    apply_font_to_run(run_sub, FONT_BODY_EN, FONT_BODY_KO, 10, COLOR_GREEN, bold=True)
    
    # Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(9.5), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = Inches(0)
    tf_title.margin_top = Inches(0)
    p_title = tf_title.paragraphs[0]
    run_title = p_title.add_run()
    run_title.text = title_text
    apply_font_to_run(run_title, FONT_TITLE_EN, FONT_TITLE_KO, 26, COLOR_DARK_TEXT)
    
    # Page Number
    page_box = slide.shapes.add_textbox(Inches(11.533), Inches(0.5), Inches(1.0), Inches(0.4))
    tf_page = page_box.text_frame
    tf_page.margin_right = Inches(0)
    tf_page.margin_top = Inches(0)
    p_page = tf_page.paragraphs[0]
    p_page.alignment = PP_ALIGN.RIGHT
    run_page = p_page.add_run()
    run_page.text = page_num
    apply_font_to_run(run_page, FONT_BODY_EN, FONT_BODY_KO, 11, COLOR_GREEN, bold=True)

def create_card_with_text_box(slide, left, top, width, height, bg_color, border_color=None, border_width_pt=1.5):
    """Creates a rounded rectangle card and returns a text frame placed exactly inside it."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(border_width_pt)
    else:
        card.line.fill.background()
    
    # Add a separate transparent text box inside the card for perfect formatting
    margin = Inches(0.2)
    tb = slide.shapes.add_textbox(left + margin, top + margin, width - (margin * 2), height - (margin * 2))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    return tf

def add_runs_to_text_frame(tf, run_specs):
    """Populates a text frame with formatted paragraphs and multiple runs based on specifications."""
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = ""
    
    for i, spec in enumerate(run_specs):
        if i == 0:
            p = p0
        else:
            p = tf.add_paragraph()
            
        if 'align' in spec:
            p.alignment = spec['align']
        if 'space_after' in spec:
            p.space_after = Pt(spec['space_after'])
        if 'space_before' in spec:
            p.space_before = Pt(spec['space_before'])
        if 'line_spacing' in spec:
            p.line_spacing = spec['line_spacing']
            
        # Support multiple sub-runs within a single paragraph
        if 'runs' in spec:
            for sub_spec in spec['runs']:
                run = p.add_run()
                run.text = sub_spec.get('text', '')
                apply_font_to_run(
                    run, 
                    sub_spec.get('font_en', FONT_BODY_EN), 
                    sub_spec.get('font_ko', FONT_BODY_KO), 
                    sub_spec.get('size', 12), 
                    sub_spec.get('color', COLOR_DARK_TEXT), 
                    sub_spec.get('bold', False)
                )
        else:
            run = p.add_run()
            run.text = spec.get('text', '')
            apply_font_to_run(
                run, 
                spec.get('font_en', FONT_BODY_EN), 
                spec.get('font_ko', FONT_BODY_KO), 
                spec.get('size', 12), 
                spec.get('color', COLOR_DARK_TEXT), 
                spec.get('bold', False)
            )

def set_slide_background_white(slide):
    """Sets the slide background to solid white."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE

# ---------------------------------------------------------
# PRESENTATION GENERATION
# ---------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # ---------------------------------------------------------
    # SLIDE 1: Title (표지)
    # ---------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide1)
    
    # Geometric Decoration (Overlapping Shapes in Green & Soft Tones)
    s1_shape1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(1.0), Inches(4.5), Inches(5.5))
    s1_shape1.fill.solid()
    s1_shape1.fill.fore_color.rgb = COLOR_LIGHT_GREEN
    s1_shape1.line.fill.background()
    
    s1_shape2 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.8), Inches(2.2), Inches(3.6), Inches(3.6))
    s1_shape2.fill.solid()
    s1_shape2.fill.fore_color.rgb = COLOR_GREEN
    s1_shape2.line.fill.background()
    
    s1_shape3 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(3.6), Inches(2.8), Inches(2.8))
    s1_shape3.fill.solid()
    s1_shape3.fill.fore_color.rgb = COLOR_YELLOW
    s1_shape3.line.fill.background()
    
    # Header pitch-deck tag
    tag_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.0), Inches(0.5))
    tf_tag = tag_box.text_frame
    run_tag = tf_tag.paragraphs[0].add_run()
    run_tag.text = "SERVICE PROPOSAL  |  PITCH DECK"
    apply_font_to_run(run_tag, FONT_BODY_BOLD_EN, FONT_BODY_KO_BOLD, 11, COLOR_GREEN)
    
    # Title Text Box
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(6.5), Inches(3.5))
    tf1 = title_box.text_frame
    runs_1 = [
        {
            'runs': [
                {'text': "SyncTrip", 'font_en': FONT_TITLE_EN, 'font_ko': FONT_TITLE_KO, 'size': 58, 'color': COLOR_GREEN, 'bold': True}
            ],
            'space_after': 15
        },
        {
            'runs': [
                {'text': "나의 여행 스타일과 100% 싱크되는\n신뢰받는 버디 찾기 서비스", 'font_en': FONT_BODY_EN, 'font_ko': FONT_BODY_KO, 'size': 20, 'color': COLOR_DARK_TEXT}
            ]
        }
    ]
    add_runs_to_text_frame(tf1, runs_1)
    
    # ---------------------------------------------------------
    # SLIDE 2: Service Introduction (서비스 소개)
    # ---------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide2)
    add_header(slide2, "SyncTrip | 기획 배경 및 해결 방안", "서비스 소개", "02")
    
    # Divider line
    divider = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.6), Inches(2.0), Inches(0.01), Inches(4.5))
    divider.fill.solid()
    divider.fill.fore_color.rgb = COLOR_BORDER_GRAY
    divider.line.fill.background()
    
    # Column Header: Problem
    prob_lbl_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(5.0), Inches(0.4))
    run_plbl = prob_lbl_box.text_frame.paragraphs[0].add_run()
    run_plbl.text = "기획 배경 (Problem)"
    apply_font_to_run(run_plbl, FONT_BODY_BOLD_EN, FONT_BODY_KO_BOLD, 15, COLOR_MUTED_TEXT)
    
    # Column Header: Solution
    sol_lbl_box = slide2.shapes.add_textbox(Inches(6.9), Inches(1.7), Inches(5.0), Inches(0.4))
    run_slbl = sol_lbl_box.text_frame.paragraphs[0].add_run()
    run_slbl.text = "해결 방안 (Solution)"
    apply_font_to_run(run_slbl, FONT_BODY_BOLD_EN, FONT_BODY_KO_BOLD, 15, COLOR_GREEN)
    
    # Left: 2 Problem Cards (White background, soft border, dark accents)
    prob_y1 = Inches(2.2)
    prob_y2 = Inches(4.4)
    prob_w = Inches(5.5)
    prob_h = Inches(2.0)
    
    # Problem 1 Card
    tf_p1 = create_card_with_text_box(slide2, Inches(0.8), prob_y1, prob_w, prob_h, COLOR_WHITE, COLOR_BORDER_GRAY)
    p1_runs = [
        {
            'runs': [
                {'text': "LIMIT 01  |  ", 'font_en': FONT_BODY_BOLD_EN, 'font_ko': FONT_BODY_KO_BOLD, 'size': 11, 'color': COLOR_MUTED_TEXT},
                {'text': "가치관 차이로 인한 동행 갈등", 'font_en': FONT_BODY_BOLD_EN, 'font_ko': FONT_BODY_KO_BOLD, 'size': 14, 'color': COLOR_DARK_TEXT}
            ],
            'space_after': 6
        },
        {
            'runs': [
                {'text': "• 성별, 연령대 등의 1차원적 조건에만 의존하여 실제 여행 중 식사, 쇼핑, 활동량, 가치관 차이로 심각한 트러블이 빈번하게 발생함.", 'font_en': FONT_BODY_REG_EN, 'font_ko': FONT_BODY_KO_REG, 'size': 11, 'color': COLOR_MUTED_TEXT}
            ]
        }
    ]
    add_runs_to_text_frame(tf_p1, p1_runs)
    
    # Problem 2 Card
    tf_p2 = create_card_with_text_box(slide2, Inches(0.8), prob_y2, prob_w, prob_h, COLOR_WHITE, COLOR_BORDER_GRAY)
    p2_runs = [
        {
            'runs': [
                {'text': "LIMIT 02  |  ", 'font_en': FONT_BODY_BOLD_EN, 'font_ko': FONT_BODY_KO_BOLD, 'size': 11, 'color': COLOR_MUTED_TEXT},
                {'text': "신원 보증 한계와 안전성 문제", 'font_en': FONT_BODY_BOLD_EN, 'font_ko': FONT_BODY_KO_BOLD, 'size': 14, 'color': COLOR_DARK_TEXT}
            ],
            'space_after': 6
        },
        {
            'runs': [
                {'text': "• 온라인 동행 구인의 특성상 상대방의 실제 신원을 보증하기 어려워, 노쇼(No-show) 트러블이나 범죄 노출 등 안전상의 위협 요소 존재.", 'font_en': FONT_BODY_REG_EN, 'font_ko': FONT_BODY_KO_REG, 'size': 11, 'color': COLOR_MUTED_TEXT}
            ]
        }
    ]
    add_runs_to_text_frame(tf_p2, p2_runs)
    
    # Right: 3 Solution Cards (Soft green background, vibrant text)
    sol_y1 = Inches(2.2)
    sol_y2 = Inches(3.68)
    sol_y3 = Inches(5.16)
    sol_w = Inches(5.6)
    sol_h = Inches(1.32)
    
    # Solution 1 Card
    tf_s1 = create_card_with_text_box(slide2, Inches(6.9), sol_y1, sol_w, sol_h, COLOR_LIGHT_GREEN)
    s1_runs = [
        {
            'runs': [
                {'text': "✓  정밀 성향 스코어 매칭", 'font_en': FONT_BODY_BOLD_EN, 'font_ko': FONT_BODY_KO_BOLD, 'size': 14, 'color': COLOR_GREEN}
            ],
            'space_after': 4
        },
        {
            'runs': [
                {'text': "MBTI와 12개 여행 세부 성향(음주, 흡연, 숙소 취향 등) 매칭을 통한 다차원 호환성 점수 산출 및 매칭 추천.", 'font_en': FONT_BODY_REG_EN, 'font_ko': FONT_BODY_KO_REG, 'size': 10.5, 'color': COLOR_DARK_TEXT}
            ]
        }
    ]
    add_runs_to_text_frame(tf_s1, s1_runs)
    
    # Solution 2 Card
    tf_s2 = create_card_with_text_box(slide2, Inches(6.9), sol_y2, sol_w, sol_h, COLOR_LIGHT_GREEN)
    s2_runs = [
        {
            'runs': [
                {'text': "✓  이중 안전 인증 배지 시스템", 'font_en': FONT_BODY_BOLD_EN, 'font_ko': FONT_BODY_KO_BOLD, 'size': 14, 'color': COLOR_GREEN}
            ],
            'space_after': 4
        },
        {
            'runs': [
                {'text': "카카오톡 실명 본인인증 및 소속 대학/직장 기관 이메일 인증을 도입하여 투명하고 안심할 수 있는 신뢰 체계 구축.", 'font_en': FONT_BODY_REG_EN, 'font_ko': FONT_BODY_KO_REG, 'size': 10.5, 'color': COLOR_DARK_TEXT}
            ]
        }
    ]
    add_runs_to_text_frame(tf_s2, s2_runs)
    
    # Solution 3 Card
    tf_s3 = create_card_with_text_box(slide2, Inches(6.9), sol_y3, sol_w, sol_h, COLOR_LIGHT_GREEN)
    s3_runs = [
        {
            'runs': [
                {'text': "✓  AI 여행자 페르소나 리포트", 'font_en': FONT_BODY_BOLD_EN, 'font_ko': FONT_BODY_KO_BOLD, 'size': 14, 'color': COLOR_GREEN}
            ],
            'space_after': 4
        },
        {
            'runs': [
                {'text': "설문 조사를 통해 AI가 사용자의 여행 스타일을 종합 분석하여 직관적인 칭호와 상세한 행동 요령 리포트를 제공.", 'font_en': FONT_BODY_REG_EN, 'font_ko': FONT_BODY_KO_REG, 'size': 10.5, 'color': COLOR_DARK_TEXT}
            ]
        }
    ]
    add_runs_to_text_frame(tf_s3, s3_runs)
    
    # ---------------------------------------------------------
    # SLIDE 3: Implementation Scope (서비스 구현 범위)
    # ---------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide3)
    add_header(slide3, "SyncTrip | 핵심 기능 범위", "서비스 구현 범위", "03")
    
    # Left: 4 Features Stack (White cards with soft borders)
    sc_x = Inches(0.8)
    sc_w = Inches(6.0)
    sc_h = Inches(1.15)
    sc_ys = [Inches(1.9), Inches(3.15), Inches(4.4), Inches(5.65)]
    
    features = [
        (
            "01. 동적 온보딩 및 프로필 수정",
            "MBTI, 언어, 여행 상태 및 12종의 심층 여행 성향 설정 지원. 터치형 국가 선택 슬라이드업 바텀 시트를 제공하여 사용자 편의성 극대화."
        ),
        (
            "02. 카테고리별 매칭 대시보드",
            "3:4 비율 모바일 가로 스와이프 피드 제공. '같은 MBTI 동행', '여행지 공통', '활동량 유사' 등 필터링을 통해 최적화된 동행 추천."
        ),
        (
            "03. 나비효과 성향 비교 상세",
            "상대방 상세 프로필 페이지에서 나와 일치하는 성향 요소를 일치 분석. 일치하는 성향 칩을 노란색 하이라이팅으로 표시하여 한눈에 파악."
        ),
        (
            "04. 실시간 대화 제안 및 보존",
            "앱 진입 3분 내 랜덤 사용자가 대화를 거는 실시간 드롭다운 토스트 알림 기능 구현. WebSocket & Supabase 연동으로 실시간 알림 보존."
        )
    ]
    
    for idx, (title, desc) in enumerate(features):
        tf_f = create_card_with_text_box(slide3, sc_x, sc_ys[idx], sc_w, sc_h, COLOR_WHITE, COLOR_BORDER_GRAY, border_width_pt=1.0)
        runs_f = [
            {
                'runs': [
                    {'text': title, 'font_en': FONT_BODY_BOLD_EN, 'font_ko': FONT_BODY_KO_BOLD, 'size': 12, 'color': COLOR_GREEN}
                ],
                'space_after': 3
            },
            {
                'runs': [
                    {'text': desc, 'font_en': FONT_BODY_REG_EN, 'font_ko': FONT_BODY_KO_REG, 'size': 9.5, 'color': COLOR_DARK_TEXT}
                ]
            }
        ]
        add_runs_to_text_frame(tf_f, runs_f)
        
    # Right: Smartphone Mockup to visualize mobile layout
    phone_x = Inches(7.8)
    phone_y = Inches(1.9)
    phone_w = Inches(3.6)
    phone_h = Inches(4.9)
    
    # Outer frame
    phone = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, phone_x, phone_y, phone_w, phone_h)
    phone.fill.solid()
    phone.fill.fore_color.rgb = COLOR_WHITE
    phone.line.color.rgb = COLOR_GREEN
    phone.line.width = Pt(3.5)
    
    # Speaker Notch
    notch = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, phone_x + Inches(1.2), phone_y + Inches(0.12), Inches(1.2), Inches(0.12))
    notch.fill.solid()
    notch.fill.fore_color.rgb = COLOR_GREEN
    notch.line.fill.background()
    
    # Inside Screen Background (Light soft green)
    screen = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, phone_x + Inches(0.15), phone_y + Inches(0.35), phone_w - Inches(0.3), phone_h - Inches(0.5))
    screen.fill.solid()
    screen.fill.fore_color.rgb = COLOR_LIGHT_GREEN
    screen.line.fill.background()
    
    # Mock UI: 3:4 Aspect Ratio User Card
    ui_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, phone_x + Inches(0.3), phone_y + Inches(0.5), Inches(3.0), Inches(2.8))
    ui_card.fill.solid()
    ui_card.fill.fore_color.rgb = COLOR_WHITE
    ui_card.line.color.rgb = COLOR_BORDER_GRAY
    ui_card.line.width = Pt(1)
    
    # Card Contents - User Image Circle
    pic_circ = slide3.shapes.add_shape(MSO_SHAPE.OVAL, phone_x + Inches(0.5), phone_y + Inches(0.7), Inches(0.9), Inches(0.9))
    pic_circ.fill.solid()
    pic_circ.fill.fore_color.rgb = COLOR_LIGHT_GREEN
    pic_circ.line.fill.background()
    
    # Card Contents - User Badge
    badge_rect = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, phone_x + Inches(1.6), phone_y + Inches(0.75), Inches(1.5), Inches(0.3))
    badge_rect.fill.solid()
    badge_rect.fill.fore_color.rgb = COLOR_GREEN
    badge_rect.line.fill.background()
    tf_b = badge_rect.text_frame
    tf_b.margin_left = tf_b.margin_right = tf_b.margin_top = tf_b.margin_bottom = Inches(0.01)
    r_b = tf_b.paragraphs[0].add_run()
    r_b.text = "✓ 이중 안전 배지"
    apply_font_to_run(r_b, FONT_BODY_BOLD_EN, FONT_BODY_KO_BOLD, 8, COLOR_WHITE)
    
    # Card Contents - Username Text
    name_box = slide3.shapes.add_textbox(phone_x + Inches(1.6), phone_y + Inches(1.1), Inches(1.5), Inches(0.4))
    tf_n = name_box.text_frame
    tf_n.margin_left = tf_n.margin_top = Inches(0)
    r_n = tf_n.paragraphs[0].add_run()
    r_n.text = "민수  (24)"
    apply_font_to_run(r_n, FONT_BODY_BOLD_EN, FONT_BODY_KO_BOLD, 12, COLOR_DARK_TEXT)
    
    # Card Contents - Description Line 1 & 2
    desc_l1 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, phone_x + Inches(0.5), phone_y + Inches(1.8), Inches(2.6), Inches(0.08))
    desc_l1.fill.solid()
    desc_l1.fill.fore_color.rgb = COLOR_BORDER_GRAY
    desc_l1.line.fill.background()
    
    desc_l2 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, phone_x + Inches(0.5), phone_y + Inches(1.95), Inches(2.0), Inches(0.08))
    desc_l2.fill.solid()
    desc_l2.fill.fore_color.rgb = COLOR_BORDER_GRAY
    desc_l2.line.fill.background()
    
    # Card Contents - Yellow Overlap Highlight Chips (나비효과 성향 비교)
    chip1 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, phone_x + Inches(0.5), phone_y + Inches(2.15), Inches(1.2), Inches(0.35))
    chip1.fill.solid()
    chip1.fill.fore_color.rgb = COLOR_YELLOW
    chip1.line.fill.background()
    tf_c1 = chip1.text_frame
    tf_c1.margin_left = tf_c1.margin_right = tf_c1.margin_top = tf_c1.margin_bottom = Inches(0.01)
    tf_c1.paragraphs[0].alignment = PP_ALIGN.CENTER
    r_c1 = tf_c1.paragraphs[0].add_run()
    r_c1.text = "✓ 쇼핑 선호"
    apply_font_to_run(r_c1, FONT_BODY_BOLD_EN, FONT_BODY_KO_BOLD, 8, COLOR_DARK_TEXT)
    
    chip2 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, phone_x + Inches(1.8), phone_y + Inches(2.15), Inches(1.2), Inches(0.35))
    chip2.fill.solid()
    chip2.fill.fore_color.rgb = COLOR_YELLOW
    chip2.line.fill.background()
    tf_c2 = chip2.text_frame
    tf_c2.margin_left = tf_c2.margin_right = tf_c2.margin_top = tf_c2.margin_bottom = Inches(0.01)
    tf_c2.paragraphs[0].alignment = PP_ALIGN.CENTER
    r_c2 = tf_c2.paragraphs[0].add_run()
    r_c2.text = "✓ 계획형"
    apply_font_to_run(r_c2, FONT_BODY_BOLD_EN, FONT_BODY_KO_BOLD, 8, COLOR_DARK_TEXT)
    
    # Mock UI: Real-time dropdown toast notification
    toast = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, phone_x + Inches(0.3), phone_y + Inches(3.6), Inches(3.0), Inches(0.85))
    toast.fill.solid()
    toast.fill.fore_color.rgb = COLOR_GREEN
    toast.line.fill.background()
    
    tf_t = toast.text_frame
    tf_t.margin_left = tf_t.margin_right = Inches(0.12)
    tf_t.margin_top = tf_t.margin_bottom = Inches(0.05)
    
    toast_runs = [
        {
            'runs': [
                {'text': "실시간 대화 제안 도착  |  3분 전", 'font_en': FONT_BODY_BOLD_EN, 'font_ko': FONT_BODY_KO_BOLD, 'size': 8, 'color': COLOR_LIGHT_GREEN}
            ],
            'space_after': 1
        },
        {
            'runs': [
                {'text': "'준수'님이 대화를 제안했습니다. [대화하기]", 'font_en': FONT_BODY_REG_EN, 'font_ko': FONT_BODY_KO_REG, 'size': 9, 'color': COLOR_WHITE}
            ]
        }
    ]
    add_runs_to_text_frame(tf_t, toast_runs)
    
    # ---------------------------------------------------------
    # SLIDE 4: UI & Work Process (작업 프로세스)
    # ---------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide4)
    add_header(slide4, "SyncTrip | 작업 프로세스 및 UI 설계", "서비스 UI 및 레퍼런스", "04")
    
    # Process Map: 5 steps Flow Diagram (도식화)
    step_y = Inches(1.9)
    step_w = Inches(2.1)
    step_h = Inches(1.1)
    step_xs = [Inches(0.8), Inches(3.22), Inches(5.64), Inches(8.06), Inches(10.48)]
    
    steps = [
        ("STEP 01", "온보딩\n(성향 조사)"),
        ("STEP 02", "매칭 피드\n탐색"),
        ("STEP 03", "상대 상세비교\n(나비효과)"),
        ("STEP 04", "매칭 신청/\n대화 제안"),
        ("STEP 05", "1:1 채팅\n소통")
    ]
    
    # Connectors (Triangles pointing right)
    for idx in range(len(step_xs) - 1):
        conn_x = step_xs[idx] + step_w + Inches(0.06)
        conn_y = step_y + step_h / 2 - Inches(0.15)
        conn = slide4.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, conn_x, conn_y, Inches(0.2), Inches(0.3))
        conn.rotation = 90
        conn.fill.solid()
        conn.fill.fore_color.rgb = COLOR_GREEN
        conn.line.fill.background()
        
    # Draw Process Step Nodes
    for idx, x_pos in enumerate(step_xs):
        tf_step = create_card_with_text_box(slide4, x_pos, step_y, step_w, step_h, COLOR_LIGHT_GREEN)
        tf_step.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        runs_step = [
            {
                'runs': [
                    {'text': steps[idx][0], 'font_en': FONT_BODY_BOLD_EN, 'font_ko': FONT_BODY_KO_BOLD, 'size': 9, 'color': COLOR_GREEN}
                ],
                'space_after': 4,
                'align': PP_ALIGN.CENTER
            },
            {
                'runs': [
                    {'text': steps[idx][1], 'font_en': FONT_BODY_BOLD_EN, 'font_ko': FONT_BODY_KO_BOLD, 'size': 11.5, 'color': COLOR_DARK_TEXT}
                ],
                'align': PP_ALIGN.CENTER
            }
        ]
        add_runs_to_text_frame(tf_step, runs_step)
        
    # Design Identity Cards (3 side-by-side)
    id_y = Inches(3.5)
    id_w = Inches(3.64)
    id_h = Inches(3.2)
    id_xs = [Inches(0.8), Inches(4.84), Inches(8.88)]
    
    id_contents = [
        (
            "네이버/라인 테마 컬러 일체화",
            "메인 포인트 컬러로 네이버/라인 테마 색상(#00C73C / #00BC70)을 선정하여 신뢰도와 역동성이 가미된 현대적 분위기를 고취합니다. 깔끔하고 밝은 톤의 녹색 계열을 메인으로 사용해 젊은 사용자들의 만족도를 극대화했습니다."
        ),
        (
            "레퍼런스 디자인 100% 반영",
            "3:4 가로 스와이프 레이아웃, 신뢰할 수 있는 사용자 배지(카카오톡 본인인증 및 소속 직장/학교 인증 마크), 초록색 체크(✓) 등의 기획안 요소를 완벽히 재현하여 기획 의도를 충실하게 시각화하였습니다."
        ),
        (
            "반응형 브라우저 레이아웃",
            "모바일 화면 가로/세로 비율(Mobile Aspect Ratio)이 데스크톱 브라우저 환경에서 늘어나거나 찌그러지지 않고 좌우 마진이 잡힌 상태로 원본 비율이 유지되도록 최적화된 반응형 목업 배치를 연출하였습니다."
        )
    ]
    
    for idx, x_pos in enumerate(id_xs):
        tf_id = create_card_with_text_box(slide4, x_pos, id_y, id_w, id_h, COLOR_WHITE, COLOR_GREEN, border_width_pt=1.5)
        runs_id = [
            {
                'runs': [
                    {'text': f"✓  {id_contents[idx][0]}", 'font_en': FONT_BODY_BOLD_EN, 'font_ko': FONT_BODY_KO_BOLD, 'size': 13.5, 'color': COLOR_GREEN}
                ],
                'space_after': 10
            },
            {
                'runs': [
                    {'text': id_contents[idx][1], 'font_en': FONT_BODY_REG_EN, 'font_ko': FONT_BODY_KO_REG, 'size': 10.5, 'color': COLOR_DARK_TEXT}
                ],
                'line_spacing': 1.15
            }
        ]
        add_runs_to_text_frame(tf_id, runs_id)
        
    # ---------------------------------------------------------
    # SLIDE 5: Backend & DB Config (백엔드 및 DB 구성)
    # ---------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide5)
    add_header(slide5, "SyncTrip | 아키텍처 및 DB 스키마", "백엔드 및 DB 구성 (ERD)", "05")
    
    # Left: 3 Stacked System Cards (Soft Green fill)
    sys_x = Inches(0.8)
    sys_w = Inches(5.4)
    sys_h = Inches(1.3)
    sys_ys = [Inches(2.0), Inches(3.55), Inches(5.1)]
    
    systems = [
        (
            "클라우드 Supabase DB 연동",
            "PostgreSQL 기반의 Supabase 클라우드 데이터베이스 인프라 구축을 신속하게 완료하고, 실시간 연결 서비스를 바탕으로 유기적인 통신 구조를 실현하였습니다."
        ),
        (
            "실시간 실소통 동기화 처리",
            "클라이언트 동작과 Supabase DB 테이블 간의 양방향 연동을 통한 실시간 매칭 동기화 처리를 구축하여 끊김 없고 부드러운 대화 제안 프로세스를 보장합니다."
        ),
        (
            "예외 처리 (Fallback 하이브리드)",
            "Supabase DB Key 미설정 및 오프라인 상태에서도 LocalStorage 데이터 엔진으로 원활히 작동을 이어갈 수 있도록 보장하는 안정적 아키텍처를 적용했습니다."
        )
    ]
    
    for idx, y_pos in enumerate(sys_ys):
        tf_sys = create_card_with_text_box(slide5, sys_x, y_pos, sys_w, sys_h, COLOR_LIGHT_GREEN)
        runs_sys = [
            {
                'runs': [
                    {'text': f"▪  {systems[idx][0]}", 'font_en': FONT_BODY_BOLD_EN, 'font_ko': FONT_BODY_KO_BOLD, 'size': 13.5, 'color': COLOR_GREEN}
                ],
                'space_after': 6
            },
            {
                'runs': [
                    {'text': systems[idx][1], 'font_en': FONT_BODY_REG_EN, 'font_ko': FONT_BODY_KO_REG, 'size': 10.5, 'color': COLOR_DARK_TEXT}
                ],
                'line_spacing': 1.1
            }
        ]
        add_runs_to_text_frame(tf_sys, runs_sys)
        
    # Right: ERD Diagram Picture Frame
    erd_frame_x = Inches(6.6)
    erd_frame_y = Inches(2.0)
    erd_frame_w = Inches(5.933)
    erd_frame_h = Inches(4.4)
    
    # Outer Frame
    erd_frame = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, erd_frame_x, erd_frame_y, erd_frame_w, erd_frame_h)
    erd_frame.fill.solid()
    erd_frame.fill.fore_color.rgb = COLOR_WHITE
    erd_frame.line.color.rgb = COLOR_GREEN
    erd_frame.line.width = Pt(1.5)
    
    # Insert ERD Image if exists
    if os.path.exists(ERD_IMAGE_PATH):
        margin_img = Inches(0.12)
        slide5.shapes.add_picture(
            ERD_IMAGE_PATH, 
            erd_frame_x + margin_img, 
            erd_frame_y + margin_img, 
            erd_frame_w - (margin_img * 2), 
            erd_frame_h - (margin_img * 2)
        )
    else:
        # Fallback Text
        fallback_tb = slide5.shapes.add_textbox(erd_frame_x, erd_frame_y + Inches(1.8), erd_frame_w, Inches(1.0))
        tf_fb = fallback_tb.text_frame
        tf_fb.paragraphs[0].alignment = PP_ALIGN.CENTER
        run_fb = tf_fb.paragraphs[0].add_run()
        run_fb.text = "[ ERD Diagram Image Missing ]\nPlease place synctrip_clean_erd_diagram.jpg in the workspace."
        apply_font_to_run(run_fb, FONT_BODY_BOLD_EN, FONT_BODY_KO_BOLD, 12, COLOR_GREEN)
        
    # ---------------------------------------------------------
    # SLIDE 6: Demo Script (서비스 이용 시연 흐름)
    # ---------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide6)
    add_header(slide6, "SyncTrip | 핵심 기능 시나리오", "서비스 이용 시연 흐름", "06")
    
    # 5 Step Horizontal Roadmap Layout
    roadmap_y = Inches(2.2)
    roadmap_w = Inches(2.18)
    roadmap_h = Inches(4.2)
    roadmap_xs = [Inches(0.8), Inches(3.2), Inches(5.6), Inches(8.0), Inches(10.4)]
    
    demo_steps = [
        (
            "1단계",
            "프로필/국가선택",
            "마이페이지에서 수정 클릭 후 신규 '국가 선택 바텀 시트' 시연. 가고 싶은 나라(일본, 몽골 등)를 직관적으로 선택하고 저장하는 UX 흐름 강조."
        ),
        (
            "2단계",
            "대시보드 추천",
            "성향 정보가 입력되면 대시보드 추천 리스트에 상대 카드 실시간 집계. 3:4 스와이프 카드를 가로로 넘기며 나와 잘 맞는 사람들의 프로필 리스트 확인."
        ),
        (
            "3단계",
            "상세 매치 분석",
            "추천된 카드 클릭 시 '나비효과 성향 비교' 상세로 이동. 상대와 나의 성향 중 겹치는 항목이 노란색 칩으로 하이라이트 표현되는 효과 시연."
        ),
        (
            "4단계",
            "대화 제안 알림",
            "대시보드를 탐색하는 중에 3분 이내로 화면 상단에서 스프링 바운스 애니메이션과 함께 드롭다운되는 대화 제안 실시간 토스트 알림 수신."
        ),
        (
            "5단계",
            "채팅방 및 보존",
            "알림 클릭 시 다이렉트로 대화방 이동. Llama3 기반 AI가 매칭 상대 성향에 맞춰 응답하는 동적 채팅과 방 복귀 시에도 유지되는 영구성 증명."
        )
    ]
    
    # Roadmap Center Connector Line
    conn_line = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), roadmap_y + Inches(0.4), Inches(11.78), Inches(0.02))
    conn_line.fill.solid()
    conn_line.fill.fore_color.rgb = COLOR_LIGHT_GREEN
    conn_line.line.fill.background()
    
    for idx, x_pos in enumerate(roadmap_xs):
        # Draw Card
        tf_demo = create_card_with_text_box(slide6, x_pos, roadmap_y + Inches(0.8), roadmap_w, roadmap_h - Inches(0.8), COLOR_WHITE, COLOR_BORDER_GRAY, border_width_pt=1.0)
        
        runs_demo = [
            {
                'runs': [
                    {'text': demo_steps[idx][1], 'font_en': FONT_BODY_BOLD_EN, 'font_ko': FONT_BODY_KO_BOLD, 'size': 12.5, 'color': COLOR_DARK_TEXT}
                ],
                'space_after': 6,
                'align': PP_ALIGN.CENTER
            },
            {
                'runs': [
                    {'text': demo_steps[idx][2], 'font_en': FONT_BODY_REG_EN, 'font_ko': FONT_BODY_KO_REG, 'size': 10, 'color': COLOR_MUTED_TEXT}
                ],
                'line_spacing': 1.15
            }
        ]
        add_runs_to_text_frame(tf_demo, runs_demo)
        
        # Step Circle badge on the line
        circ = slide6.shapes.add_shape(MSO_SHAPE.OVAL, x_pos + roadmap_w / 2 - Inches(0.35), roadmap_y + Inches(0.05), Inches(0.7), Inches(0.7))
        circ.fill.solid()
        circ.fill.fore_color.rgb = COLOR_GREEN
        circ.line.fill.background()
        tf_c = circ.text_frame
        tf_c.margin_left = tf_c.margin_right = tf_c.margin_top = tf_c.margin_bottom = Inches(0.01)
        tf_c.paragraphs[0].alignment = PP_ALIGN.CENTER
        r_c = tf_c.paragraphs[0].add_run()
        r_c.text = str(idx + 1)
        apply_font_to_run(r_c, FONT_TITLE_EN, FONT_TITLE_KO, 13, COLOR_WHITE)
        
        # Extra Design Highlights based on slide requirements
        if idx == 2: # Step 3: yellow chip demo
            y_chip = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos + Inches(0.35), roadmap_y + roadmap_h - Inches(0.45), Inches(1.48), Inches(0.32))
            y_chip.fill.solid()
            y_chip.fill.fore_color.rgb = COLOR_YELLOW
            y_chip.line.fill.background()
            tf_yc = y_chip.text_frame
            tf_yc.margin_left = tf_yc.margin_right = tf_yc.margin_top = tf_yc.margin_bottom = Inches(0.01)
            tf_yc.paragraphs[0].alignment = PP_ALIGN.CENTER
            r_yc = tf_yc.paragraphs[0].add_run()
            r_yc.text = "✓ Overlap Highlight"
            apply_font_to_run(r_yc, FONT_BODY_BOLD_EN, FONT_BODY_KO_BOLD, 7.5, COLOR_DARK_TEXT)
            
        elif idx == 3: # Step 4: Toast Popup demo
            t_chip = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos + Inches(0.35), roadmap_y + roadmap_h - Inches(0.45), Inches(1.48), Inches(0.32))
            t_chip.fill.solid()
            t_chip.fill.fore_color.rgb = COLOR_GREEN
            t_chip.line.fill.background()
            tf_tc = t_chip.text_frame
            tf_tc.margin_left = tf_tc.margin_right = tf_tc.margin_top = tf_tc.margin_bottom = Inches(0.01)
            tf_tc.paragraphs[0].alignment = PP_ALIGN.CENTER
            r_tc = tf_tc.paragraphs[0].add_run()
            r_tc.text = "✉ 대화 알림 수신"
            apply_font_to_run(r_tc, FONT_BODY_BOLD_EN, FONT_BODY_KO_BOLD, 7.5, COLOR_WHITE)
            
        elif idx == 4: # Step 5: Chat bubble demo
            c_chip = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos + Inches(0.35), roadmap_y + roadmap_h - Inches(0.45), Inches(1.48), Inches(0.32))
            c_chip.fill.solid()
            c_chip.fill.fore_color.rgb = COLOR_LIGHT_GREEN
            c_chip.line.fill.background()
            tf_cc = c_chip.text_frame
            tf_cc.margin_left = tf_cc.margin_right = tf_cc.margin_top = tf_cc.margin_bottom = Inches(0.01)
            tf_cc.paragraphs[0].alignment = PP_ALIGN.CENTER
            r_cc = tf_cc.paragraphs[0].add_run()
            r_cc.text = "💬 실시간 소통"
            apply_font_to_run(r_cc, FONT_BODY_BOLD_EN, FONT_BODY_KO_BOLD, 7.5, COLOR_DARK_TEXT)
            
    # ---------------------------------------------------------
    # SLIDE 7: End (감사합니다)
    # ---------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide7)
    
    # Overlapping circle decos
    end_shape1 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(1.0), Inches(4.5), Inches(5.5))
    end_shape1.fill.solid()
    end_shape1.fill.fore_color.rgb = COLOR_LIGHT_GREEN
    end_shape1.line.fill.background()
    
    end_shape2 = slide7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.8), Inches(2.2), Inches(3.6), Inches(3.6))
    end_shape2.fill.solid()
    end_shape2.fill.fore_color.rgb = COLOR_GREEN
    end_shape2.line.fill.background()
    
    end_shape3 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(3.6), Inches(2.8), Inches(2.8))
    end_shape3.fill.solid()
    end_shape3.fill.fore_color.rgb = COLOR_YELLOW
    end_shape3.line.fill.background()
    
    # Ending Text Box
    end_box = slide7.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(6.5), Inches(3.0))
    tf_end = end_box.text_frame
    runs_end = [
        {
            'runs': [
                {'text': "감사합니다", 'font_en': FONT_TITLE_EN, 'font_ko': FONT_TITLE_KO, 'size': 58, 'color': COLOR_GREEN, 'bold': True}
            ],
            'space_after': 15
        },
        {
            'runs': [
                {'text': "SyncTrip  |  나의 여행 스타일과 100% 싱크되는\n신뢰받는 버디 찾기 서비스 소개", 'font_en': FONT_BODY_EN, 'font_ko': FONT_BODY_KO, 'size': 18, 'color': COLOR_DARK_TEXT}
            ]
        }
    ]
    add_runs_to_text_frame(tf_end, runs_end)
    
    # Save Presentation
    prs.save(OUTPUT_PATH)
    print(f"Presentation saved successfully to: {OUTPUT_PATH}")
    
    # Verify Fonts
    verify_fonts(OUTPUT_PATH)

def verify_fonts(output_path):
    """Verifies that all text runs inside the presentation are using the requested fonts."""
    prs = Presentation(output_path)
    allowed_fonts = {FONT_TITLE_EN, FONT_BODY_EN, FONT_BODY_BOLD_EN, FONT_BODY_REG_EN, FONT_BODY_LIGHT_EN}
    allowed_fonts_ko = {FONT_TITLE_KO, FONT_BODY_KO, FONT_BODY_KO_BOLD, FONT_BODY_KO_REG, FONT_BODY_KO_LIGHT}
    all_ok = True
    
    print("\n" + "="*70)
    print("FONT VERIFICATION LOG")
    print("="*70)
    
    for s_idx, slide in enumerate(prs.slides):
        print(f"\n[Slide {s_idx + 1}]")
        for sh_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                for p_idx, p in enumerate(shape.text_frame.paragraphs):
                    for r_idx, r in enumerate(p.runs):
                        font_name = r.font.name
                        
                        # Find EA typeface from XML
                        rPr = r.font._element
                        ea_font = None
                        for child in rPr:
                            if child.tag.endswith('ea'):
                                ea_font = child.get('typeface')
                        
                        is_latin_ok = font_name in allowed_fonts
                        is_ea_ok = ea_font in allowed_fonts_ko
                        
                        status = "OK"
                        if not is_latin_ok and not is_ea_ok:
                            status = "WARNING"
                            all_ok = False
                            
                        safe_text = r.text[:20].encode('ascii', errors='ignore').decode('ascii')
                        print(f"  Shape {sh_idx+1} | Text: '{safe_text}...' | Latin Font: {font_name} (ok={is_latin_ok}) | EA Font: {ea_font} (ok={is_ea_ok}) -> {status}")
                        
    print("\n" + "="*70)
    if all_ok:
        print("VERIFICATION RESULT: SUCCESS (All fonts verified)")
    else:
        print("VERIFICATION RESULT: WARNING (Some text runs missing target fonts)")
    print("="*70)

if __name__ == "__main__":
    main()
