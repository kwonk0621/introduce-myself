import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from spire.presentation import Presentation as SpirePresentation, FileFormat as SpireFileFormat

# ---------------------------------------------------------
# CONSTANTS & CONFIGURATION
# ---------------------------------------------------------
OUTPUT_PPTX_PATH = "SyncTrip_Presentation.pptx"
OUTPUT_PDF_PATH = "SyncTrip_Presentation.pdf"
ERD_IMAGE_PATH = "synctrip_clean_erd_diagram.jpg"

# Color Palette (Premium White & Green Theme)
COLOR_WHITE = RGBColor(255, 255, 255)            # #FFFFFF
COLOR_LIGHT_GREEN = RGBColor(211, 245, 221)       # #D3F5DD (Soft light green background/accent)
COLOR_GREEN = RGBColor(0, 188, 112)              # #00BC70 (Theme Primary Green)
COLOR_DARK_TEXT = RGBColor(30, 41, 35)            # #1E2923 (Deep forest dark green - premium black replacement)
COLOR_MUTED_TEXT = RGBColor(107, 128, 117)        # #6B8075 (Muted sage gray)
COLOR_YELLOW = RGBColor(255, 245, 204)            # #FFF5CC (Soft highlight yellow for overlap chips)
COLOR_BORDER_GRAY = RGBColor(210, 225, 215)       # Soft grey-green border

# Font Settings (Gmarket Sans TTF)
FONT_TITLE = "Gmarket Sans TTF Bold"
FONT_BODY_BOLD = "Gmarket Sans TTF Bold"
FONT_BODY_MEDIUM = "Gmarket Sans TTF Medium"
FONT_BODY_LIGHT = "Gmarket Sans TTF Light"

# ---------------------------------------------------------
# HELPER FUNCTIONS
# ---------------------------------------------------------
def apply_font_to_run(run, font_name, size_pt, color_rgb, bold=False):
    """Applies font formatting to a text run, forcing Latin and East Asian font matching."""
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color_rgb
    run.font.bold = bold
    
    # Force East Asian typeface matching (required for Korean rendering in PowerPoint)
    rPr = run.font._element
    for child in list(rPr):
        if child.tag.endswith('ea'):
            rPr.remove(child)
    
    ea = OxmlElement('a:ea')
    ea.set('typeface', font_name)
    rPr.append(ea)

def add_header(slide, category_text, title_text, page_num):
    """Creates a consistent premium header template on the slide."""
    # Top thin line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_GREEN
    line.line.fill.background()
    
    # Category / Sub title
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(6.0), Inches(0.4))
    tf_sub = sub_box.text_frame
    tf_sub.margin_left = Inches(0)
    tf_sub.margin_top = Inches(0)
    p_sub = tf_sub.paragraphs[0]
    run_sub = p_sub.add_run()
    run_sub.text = f"SyncTrip | {category_text}"
    apply_font_to_run(run_sub, FONT_BODY_BOLD, 10, COLOR_GREEN)
    
    # Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(9.5), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = Inches(0)
    tf_title.margin_top = Inches(0)
    p_title = tf_title.paragraphs[0]
    run_title = p_title.add_run()
    run_title.text = title_text
    apply_font_to_run(run_title, FONT_TITLE, 24, COLOR_DARK_TEXT)
    
    # Page Number
    page_box = slide.shapes.add_textbox(Inches(11.533), Inches(0.5), Inches(1.0), Inches(0.4))
    tf_page = page_box.text_frame
    tf_page.margin_right = Inches(0)
    tf_page.margin_top = Inches(0)
    p_page = tf_page.paragraphs[0]
    p_page.alignment = PP_ALIGN.RIGHT
    run_page = p_page.add_run()
    run_page.text = page_num
    apply_font_to_run(run_page, FONT_BODY_BOLD, 11, COLOR_GREEN)

def create_card_with_text_box(slide, left, top, width, height, bg_color, border_color=None, border_width_pt=1.5, rx=0.05):
    """Creates a rounded rectangle card and returns a text frame placed exactly inside it."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(border_width_pt)
    else:
        card.line.fill.background()
    
    # Add a separate transparent text box inside the card for perfect formatting
    margin = Inches(0.2)
    tb = slide.shapes.add_textbox(left + margin, top + margin, width - (margin * 2), height - (margin * 2))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    return tf

def add_runs_to_text_frame(tf, run_specs):
    """Populates a text frame with formatted paragraphs and multiple runs based on specifications."""
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = ""
    
    for i, spec in enumerate(run_specs):
        if i == 0:
            p = p0
        else:
            p = tf.add_paragraph()
            
        if 'align' in spec:
            p.alignment = spec['align']
        if 'space_after' in spec:
            p.space_after = Pt(spec['space_after'])
        if 'space_before' in spec:
            p.space_before = Pt(spec['space_before'])
        if 'line_spacing' in spec:
            p.line_spacing = spec['line_spacing']
            
        if 'runs' in spec:
            for sub_spec in spec['runs']:
                run = p.add_run()
                run.text = sub_spec.get('text', '')
                apply_font_to_run(
                    run, 
                    sub_spec.get('font', FONT_BODY_MEDIUM), 
                    sub_spec.get('size', 12), 
                    sub_spec.get('color', COLOR_DARK_TEXT), 
                    sub_spec.get('bold', False)
                )
        else:
            run = p.add_run()
            run.text = spec.get('text', '')
            apply_font_to_run(
                run, 
                spec.get('font', FONT_BODY_MEDIUM), 
                spec.get('size', 12), 
                spec.get('color', COLOR_DARK_TEXT), 
                spec.get('bold', False)
            )

def set_slide_background_white(slide):
    """Sets the slide background to solid white."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE

# ---------------------------------------------------------
# PRESENTATION GENERATION
# ---------------------------------------------------------
def generate_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # ---------------------------------------------------------
    # SLIDE 1: Title (표지)
    # ---------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide1)
    
    # Geometric Decoration (Overlapping Shapes in Green & Soft Tones)
    s1_shape1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(1.0), Inches(4.5), Inches(5.5))
    s1_shape1.fill.solid()
    s1_shape1.fill.fore_color.rgb = COLOR_LIGHT_GREEN
    s1_shape1.line.fill.background()
    
    s1_shape2 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.0), Inches(2.2), Inches(3.6), Inches(3.6))
    s1_shape2.fill.solid()
    s1_shape2.fill.fore_color.rgb = COLOR_GREEN
    s1_shape2.line.fill.background()
    
    s1_shape3 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(3.6), Inches(2.8), Inches(2.8))
    s1_shape3.fill.solid()
    s1_shape3.fill.fore_color.rgb = COLOR_YELLOW
    s1_shape3.line.fill.background()
    
    # Header tag
    tag_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(6.0), Inches(0.5))
    tf_tag = tag_box.text_frame
    run_tag = tf_tag.paragraphs[0].add_run()
    run_tag.text = "AI-POWERED TRAVEL COMPANION MATCHING SERVICE"
    apply_font_to_run(run_tag, FONT_BODY_BOLD, 10.5, COLOR_GREEN)
    
    # Title Text Box
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(6.5), Inches(3.5))
    tf1 = title_box.text_frame
    runs_1 = [
        {
            'runs': [
                {'text': "SyncTrip", 'font': FONT_TITLE, 'size': 62, 'color': COLOR_GREEN, 'bold': True}
            ],
            'space_after': 18
        },
        {
            'runs': [
                {'text': "나의 여행 스타일과 100% 싱크되는\n신뢰받는 여행 버디 찾기 서비스", 'font': FONT_BODY_MEDIUM, 'size': 20, 'color': COLOR_DARK_TEXT}
            ],
            'space_after': 14
        },
        {
            'runs': [
                {'text': "여행 성향 분석 기반의 AI 매칭 및 실시간 채팅 플랫폼", 'font': FONT_BODY_LIGHT, 'size': 13, 'color': COLOR_MUTED_TEXT}
            ]
        }
    ]
    add_runs_to_text_frame(tf1, runs_1)
    
    # ---------------------------------------------------------
    # SLIDE 2: Service Introduction (서비스 소개)
    # ---------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide2)
    add_header(slide2, "서비스 소개", "여행 스타일 맞춤형 동행 매칭 플랫폼", "02")
    
    # Concept Introduction Box (Top Card)
    concept_tf = create_card_with_text_box(slide2, Inches(0.8), Inches(1.7), Inches(11.733), Inches(1.0), COLOR_LIGHT_GREEN)
    concept_runs = [
        {
            'runs': [
                {'text': "SyncTrip은 사용자의 성별, 연령대 뿐만 아니라 세부 여행 취향(MBTI, 활동량, 선호 국가, 흡연 여부 등)을 AI로 분석하여\n가장 호환성이 높은 최적의 여행 버디를 연결하고, 안전 인증 시스템과 실시간 채팅을 지원하는 플랫폼입니다.", 'font': FONT_BODY_MEDIUM, 'size': 13.5, 'color': COLOR_DARK_TEXT}
            ],
            'align': PP_ALIGN.CENTER
        }
    ]
    add_runs_to_text_frame(concept_tf, concept_runs)
    
    # Column Dividers / Layout (Left: Problem, Right: Solution)
    card_w = Inches(5.6)
    card_h = Inches(1.7)
    
    # Left: Problems (White background with red/muted accents, soft border)
    # Problem 1
    tf_p1 = create_card_with_text_box(slide2, Inches(0.8), Inches(3.0), card_w, card_h, COLOR_WHITE, COLOR_BORDER_GRAY)
    p1_runs = [
        {
            'runs': [
                {'text': "PROBLEM 01  |  단순 조건 검색의 한계", 'font': FONT_BODY_BOLD, 'size': 13, 'color': COLOR_MUTED_TEXT}
            ],
            'space_after': 6
        },
        {
            'runs': [
                {'text': "성별, 연령대 등 단순 정보만으로 동행을 구하면 실제 여행 중 성향(음식, 쇼핑, 활동량) 차이로 인해 심각한 의견 충돌 및 갈등이 발생합니다.", 'font': FONT_BODY_LIGHT, 'size': 11, 'color': COLOR_DARK_TEXT}
            ],
            'line_spacing': 1.15
        }
    ]
    add_runs_to_text_frame(tf_p1, p1_runs)
    
    # Problem 2
    tf_p2 = create_card_with_text_box(slide2, Inches(0.8), Inches(4.9), card_w, card_h, COLOR_WHITE, COLOR_BORDER_GRAY)
    p2_runs = [
        {
            'runs': [
                {'text': "PROBLEM 02  |  신원 불확실성과 안전 위협", 'font': FONT_BODY_BOLD, 'size': 13, 'color': COLOR_MUTED_TEXT}
            ],
            'space_after': 6
        },
        {
            'runs': [
                {'text': "익명 커뮤니티 기반의 구인은 상대방의 신원을 보증하기 어렵고, 당일 잠수(노쇼)나 안전 범죄에 노출될 수 있는 리스크를 지닙니다.", 'font': FONT_BODY_LIGHT, 'size': 11, 'color': COLOR_DARK_TEXT}
            ],
            'line_spacing': 1.15
        }
    ]
    add_runs_to_text_frame(tf_p2, p2_runs)
    
    # Right: Solutions (White background with primary green accents, solid border)
    # Solution 1
    tf_s1 = create_card_with_text_box(slide2, Inches(6.933), Inches(3.0), card_w, card_h, COLOR_WHITE, COLOR_GREEN)
    s1_runs = [
        {
            'runs': [
                {'text': "✓ SOLUTION 01  |  AI 다차원 성향 스코어링", 'font': FONT_BODY_BOLD, 'size': 13, 'color': COLOR_GREEN}
            ],
            'space_after': 6
        },
        {
            'runs': [
                {'text': "사용자가 입력한 MBTI, 활동량, 흡연여부 등 12가지 항목의 성향 프로필을 정밀 매칭 알고리즘으로 분석하여 동행 호환성을 직관적으로 제시합니다.", 'font': FONT_BODY_LIGHT, 'size': 11, 'color': COLOR_DARK_TEXT}
            ],
            'line_spacing': 1.15
        }
    ]
    add_runs_to_text_frame(tf_s1, s1_runs)
    
    # Solution 2
    tf_s2 = create_card_with_text_box(slide2, Inches(6.933), Inches(4.9), card_w, card_h, COLOR_WHITE, COLOR_GREEN)
    s2_runs = [
        {
            'runs': [
                {'text': "✓ SOLUTION 02  |  이메일 신원 인증 시스템", 'font': FONT_BODY_BOLD, 'size': 13, 'color': COLOR_GREEN}
            ],
            'space_after': 6
        },
        {
            'runs': [
                {'text': "회원가입 후 직장 혹은 학교 소속 인증 메일 링크를 거치도록 하여 안전성이 확보된 유저만을 매칭 풀에 편입하고 인증 배지를 부여합니다.", 'font': FONT_BODY_LIGHT, 'size': 11, 'color': COLOR_DARK_TEXT}
            ],
            'line_spacing': 1.15
        }
    ]
    add_runs_to_text_frame(tf_s2, s2_runs)
    
    # ---------------------------------------------------------
    # SLIDE 3: Implementation Scope (서비스 구현 범위)
    # ---------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide3)
    add_header(slide3, "서비스 구현 범위", "핵심 기능 및 구현 아키텍처", "03")
    
    # 4 Feature Card Grid (Left)
    sc_x = Inches(0.8)
    sc_w = Inches(5.6)
    sc_h = Inches(1.15)
    sc_ys = [Inches(1.8), Inches(3.1), Inches(4.4), Inches(5.7)]
    
    features = [
        (
            "01. 성향 진단 및 프로필 설정",
            "MBTI, 활동량(낮음/보통/높음), 선호 국가(바텀 시트 형식 멀티 선택), 흡연 여부 등 다양한 성향 카테고리를 직관적으로 선택하는 설정 페이지를 제공합니다."
        ),
        (
            "02. AI 기반 동행 매칭 시스템",
            "나와 타인의 프로필 메타데이터를 비교분석하여 호환성 점수를 실시간으로 연산하고, 대시보드 피드를 통해 맞춤형 추천 동행 카드를 제시합니다."
        ),
        (
            "03. 나비효과 성향 비교 상세",
            "동행 카드 클릭 시 진입하는 상세 페이지에서, 나와 상대방의 여행 성향 중 겹치는 일치 항목을 노란색 하이라이팅 배지로 표현해 직관성을 보장합니다."
        ),
        (
            "04. 안전 인증 및 실시간 채팅",
            "학교/직장 이메일 본인 인증 시스템으로 가입 신뢰성을 검증하며, WebSocket & Supabase DB 트리거로 1:1 실시간 대화 및 실시간 알림 토스트를 지원합니다."
        )
    ]
    
    for idx, (title, desc) in enumerate(features):
        tf_f = create_card_with_text_box(slide3, sc_x, sc_ys[idx], sc_w, sc_h, COLOR_WHITE, COLOR_BORDER_GRAY, border_width_pt=1.0)
        runs_f = [
            {
                'runs': [
                    {'text': title, 'font': FONT_BODY_BOLD, 'size': 12, 'color': COLOR_GREEN}
                ],
                'space_after': 4
            },
            {
                'runs': [
                    {'text': desc, 'font': FONT_BODY_LIGHT, 'size': 9.5, 'color': COLOR_DARK_TEXT}
                ],
                'line_spacing': 1.1
            }
        ]
        add_runs_to_text_frame(tf_f, runs_f)
        
    # Smartphone Mockup Visual on the Right (Mockup of App UI)
    phone_x = Inches(7.8)
    phone_y = Inches(1.8)
    phone_w = Inches(3.6)
    phone_h = Inches(5.05)
    
    # Outer frame
    phone = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, phone_x, phone_y, phone_w, phone_h)
    phone.fill.solid()
    phone.fill.fore_color.rgb = COLOR_WHITE
    phone.line.color.rgb = COLOR_GREEN
    phone.line.width = Pt(3.5)
    
    # Speaker Notch
    notch = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, phone_x + Inches(1.2), phone_y + Inches(0.12), Inches(1.2), Inches(0.12))
    notch.fill.solid()
    notch.fill.fore_color.rgb = COLOR_GREEN
    notch.line.fill.background()
    
    # Inside Screen Background (Light soft green)
    screen = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, phone_x + Inches(0.15), phone_y + Inches(0.35), phone_w - Inches(0.3), phone_h - Inches(0.5))
    screen.fill.solid()
    screen.fill.fore_color.rgb = COLOR_LIGHT_GREEN
    screen.line.fill.background()
    
    # Mock UI: 3:4 Aspect Ratio User Card
    ui_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, phone_x + Inches(0.3), phone_y + Inches(0.5), Inches(3.0), Inches(2.9))
    ui_card.fill.solid()
    ui_card.fill.fore_color.rgb = COLOR_WHITE
    ui_card.line.color.rgb = COLOR_BORDER_GRAY
    ui_card.line.width = Pt(1)
    
    # Card Contents - User Image Circle
    pic_circ = slide3.shapes.add_shape(MSO_SHAPE.OVAL, phone_x + Inches(0.5), phone_y + Inches(0.7), Inches(0.9), Inches(0.9))
    pic_circ.fill.solid()
    pic_circ.fill.fore_color.rgb = COLOR_LIGHT_GREEN
    pic_circ.line.fill.background()
    
    # Card Contents - User Badge
    badge_rect = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, phone_x + Inches(1.6), phone_y + Inches(0.75), Inches(1.4), Inches(0.3))
    badge_rect.fill.solid()
    badge_rect.fill.fore_color.rgb = COLOR_GREEN
    badge_rect.line.fill.background()
    tf_b = badge_rect.text_frame
    tf_b.margin_left = tf_b.margin_right = tf_b.margin_top = tf_b.margin_bottom = Inches(0.01)
    r_b = tf_b.paragraphs[0].add_run()
    r_b.text = "✓ 이메일인증"
    apply_font_to_run(r_b, FONT_BODY_BOLD, 7.5, COLOR_WHITE)
    
    # Card Contents - Username Text
    name_box = slide3.shapes.add_textbox(phone_x + Inches(1.6), phone_y + Inches(1.15), Inches(1.5), Inches(0.4))
    tf_n = name_box.text_frame
    tf_n.margin_left = tf_n.margin_top = Inches(0)
    r_n = tf_n.paragraphs[0].add_run()
    r_n.text = "이지민  (25)"
    apply_font_to_run(r_n, FONT_BODY_BOLD, 12, COLOR_DARK_TEXT)
    
    # Card Contents - Description Line 1 & 2
    desc_l1 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, phone_x + Inches(0.5), phone_y + Inches(1.8), Inches(2.6), Inches(0.06))
    desc_l1.fill.solid()
    desc_l1.fill.fore_color.rgb = COLOR_BORDER_GRAY
    desc_l1.line.fill.background()
    
    desc_l2 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, phone_x + Inches(0.5), phone_y + Inches(1.95), Inches(2.0), Inches(0.06))
    desc_l2.fill.solid()
    desc_l2.fill.fore_color.rgb = COLOR_BORDER_GRAY
    desc_l2.line.fill.background()
    
    # Card Contents - Yellow Overlap Highlight Chips (나비효과 성향 비교)
    chip1 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, phone_x + Inches(0.5), phone_y + Inches(2.2), Inches(1.15), Inches(0.35))
    chip1.fill.solid()
    chip1.fill.fore_color.rgb = COLOR_YELLOW
    chip1.line.fill.background()
    tf_c1 = chip1.text_frame
    tf_c1.margin_left = tf_c1.margin_right = tf_c1.margin_top = tf_c1.margin_bottom = Inches(0.01)
    tf_c1.paragraphs[0].alignment = PP_ALIGN.CENTER
    r_c1 = tf_c1.paragraphs[0].add_run()
    r_c1.text = "✓ 계획형 (J)"
    apply_font_to_run(r_c1, FONT_BODY_BOLD, 7.5, COLOR_DARK_TEXT)
    
    chip2 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, phone_x + Inches(1.75), phone_y + Inches(2.2), Inches(1.15), Inches(0.35))
    chip2.fill.solid()
    chip2.fill.fore_color.rgb = COLOR_YELLOW
    chip2.line.fill.background()
    tf_c2 = chip2.text_frame
    tf_c2.margin_left = tf_c2.margin_right = tf_c2.margin_top = tf_c2.margin_bottom = Inches(0.01)
    tf_c2.paragraphs[0].alignment = PP_ALIGN.CENTER
    r_c2 = tf_c2.paragraphs[0].add_run()
    r_c2.text = "✓ 활동량 높음"
    apply_font_to_run(r_c2, FONT_BODY_BOLD, 7.5, COLOR_DARK_TEXT)
    
    # Mock UI: Real-time Toast Dropdown
    toast = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, phone_x + Inches(0.3), phone_y + Inches(3.6), Inches(3.0), Inches(0.85))
    toast.fill.solid()
    toast.fill.fore_color.rgb = COLOR_GREEN
    toast.line.fill.background()
    
    tf_t = toast.text_frame
    tf_t.margin_left = tf_t.margin_right = Inches(0.12)
    tf_t.margin_top = tf_t.margin_bottom = Inches(0.05)
    
    toast_runs = [
        {
            'runs': [
                {'text': "실시간 대화 제안 도착  |  방금 전", 'font': FONT_BODY_BOLD, 'size': 8, 'color': COLOR_LIGHT_GREEN}
            ],
            'space_after': 1
        },
        {
            'runs': [
                {'text': "'김민수'님이 대화방 참여를 요청하셨습니다.", 'font': FONT_BODY_MEDIUM, 'size': 8.5, 'color': COLOR_WHITE}
            ]
        }
    ]
    add_runs_to_text_frame(tf_t, toast_runs)
    
    # ---------------------------------------------------------
    # SLIDE 4: Service UI & Reference (작업 프로세스)
    # ---------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide4)
    add_header(slide4, "서비스 UI 및 작업 프로세스", "프로세스 플로우 및 디자인 아이덴티티", "04")
    
    # Process Map (Top Section: 5 Step Cards Flow)
    step_y = Inches(1.8)
    step_w = Inches(2.1)
    step_h = Inches(1.1)
    step_xs = [Inches(0.8), Inches(3.22), Inches(5.64), Inches(8.06), Inches(10.48)]
    
    steps = [
        ("STEP 01", "온보딩\n(성향 진단)"),
        ("STEP 02", "매칭 피드\n탐색 (3:4 카드)"),
        ("STEP 03", "상대 상세비교\n(나비효과)"),
        ("STEP 04", "매칭 신청/\n대화 제안"),
        ("STEP 05", "1:1 실시간\n채팅 소통")
    ]
    
    # Process Connector Lines (arrows)
    for idx in range(len(step_xs) - 1):
        conn_x = step_xs[idx] + step_w + Inches(0.06)
        conn_y = step_y + step_h / 2 - Inches(0.15)
        conn = slide4.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, conn_x, conn_y, Inches(0.2), Inches(0.3))
        conn.rotation = 90
        conn.fill.solid()
        conn.fill.fore_color.rgb = COLOR_GREEN
        conn.line.fill.background()
        
    for idx, x_pos in enumerate(step_xs):
        tf_step = create_card_with_text_box(slide4, x_pos, step_y, step_w, step_h, COLOR_LIGHT_GREEN)
        tf_step.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        runs_step = [
            {
                'runs': [
                    {'text': steps[idx][0], 'font': FONT_BODY_BOLD, 'size': 9, 'color': COLOR_GREEN}
                ],
                'space_after': 4,
                'align': PP_ALIGN.CENTER
            },
            {
                'runs': [
                    {'text': steps[idx][1], 'font': FONT_BODY_BOLD, 'size': 11, 'color': COLOR_DARK_TEXT}
                ],
                'align': PP_ALIGN.CENTER
            }
        ]
        add_runs_to_text_frame(tf_step, runs_step)
        
    # Design Identity (Bottom Section: 3 Columns)
    id_y = Inches(3.3)
    id_w = Inches(3.64)
    id_h = Inches(3.4)
    id_xs = [Inches(0.8), Inches(4.84), Inches(8.88)]
    
    id_contents = [
        (
            "네이버/라인 테마 컬러 일체화",
            "• 포인트 색상으로 네이버 및 라인의 브랜드 그린 테마 컬러인 #00C73C 및 #00BC70을 차용하여 활기차고 직관적인 UI 디자인 감성을 극대화합니다.\n• 눈의 피로도를 낮추고 모던한 느낌을 전달하는 화이트 & 그린 톤을 화면 전체에 조화롭게 연출하였습니다."
        ),
        (
            "레퍼런스 디자인의 완벽한 수용",
            "• 피드 화면 내 카드 레이아웃의 최적 비율인 3:4 황금비를 정밀 구현하여 시인성을 높였습니다.\n• 신뢰를 상징하는 이중 안전 인증 배지 시스템 및 초록색 체크(✓) 아이콘 등 기획 문서에서 명시된 디자인 요소들을 화면 전반에 일치시켰습니다."
        ),
        (
            "반응형 브라우저 정밀 레이아웃",
            "• PC와 모바일 디바이스 환경 전반에서 끊김 없는 사용자 경험을 보장합니다.\n• 데스크톱 브라우저 환경에서도 모바일 화면비(Mobile Aspect Ratio)가 가로로 찌그러지거나 왜곡되지 않고 좌우 여백 속에서 정밀하게 고정되도록 디자인 아키텍처를 수립했습니다."
        )
    ]
    
    for idx, x_pos in enumerate(id_xs):
        tf_id = create_card_with_text_box(slide4, x_pos, id_y, id_w, id_h, COLOR_WHITE, COLOR_GREEN, border_width_pt=1.5)
        runs_id = [
            {
                'runs': [
                    {'text': f"✓  {id_contents[idx][0]}", 'font': FONT_BODY_BOLD, 'size': 13.5, 'color': COLOR_GREEN}
                ],
                'space_after': 10
            },
            {
                'runs': [
                    {'text': id_contents[idx][1], 'font': FONT_BODY_LIGHT, 'size': 10, 'color': COLOR_DARK_TEXT}
                ],
                'line_spacing': 1.2
            }
        ]
        add_runs_to_text_frame(tf_id, runs_id)
        
    # ---------------------------------------------------------
    # SLIDE 5: Backend & DB Config (백엔드 및 DB 구성)
    # ---------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide5)
    add_header(slide5, "백엔드 및 DB 구성", "Supabase 데이터베이스 및 ERD 아키텍처", "05")
    
    # Left: 3 Stacked System Description Cards
    sys_x = Inches(0.8)
    sys_w = Inches(5.4)
    sys_h = Inches(1.4)
    sys_ys = [Inches(1.8), Inches(3.4), Inches(5.0)]
    
    systems = [
        (
            "Supabase 클라우드 데이터베이스 연동",
            "PostgreSQL 기반의 Supabase 클라우드 인프라를 활용하여 사용자 성향 데이터(MBTI, 흡연, 선호국가 등) 및 실시간 채팅 대화방 테이블을 체계적으로 구조화하고 실시간 통신을 보장합니다."
        ),
        (
            "실시간 매칭 및 동기화 처리",
            "Supabase Database trigger와 WebSocket 기술을 연동하여 지연이 존재하지 않는 매칭 생성 프로세스 및 1:1 채팅방 메시지 실시간 읽기/쓰기를 유기적으로 연동하였습니다."
        ),
        (
            "오프라인 하이브리드 Fallback 설계",
            "Supabase DB Key 미연결이나 일시적인 네트워크 오프라인 조건에서도 정상적인 시연 및 가동이 이어질 수 있도록 로컬 스토리지(LocalStorage) 기반 Fallback 엔진을 병행 탑재하였습니다."
        )
    ]
    
    for idx, y_pos in enumerate(sys_ys):
        tf_sys = create_card_with_text_box(slide5, sys_x, y_pos, sys_w, sys_h, COLOR_LIGHT_GREEN)
        runs_sys = [
            {
                'runs': [
                    {'text': f"▪  {systems[idx][0]}", 'font': FONT_BODY_BOLD, 'size': 13.5, 'color': COLOR_GREEN}
                ],
                'space_after': 6
            },
            {
                'runs': [
                    {'text': systems[idx][1], 'font': FONT_BODY_LIGHT, 'size': 10, 'color': COLOR_DARK_TEXT}
                ],
                'line_spacing': 1.15
            }
        ]
        add_runs_to_text_frame(tf_sys, runs_sys)
        
    # Right: ERD Diagram Picture Frame
    erd_frame_x = Inches(6.6)
    erd_frame_y = Inches(1.8)
    erd_frame_w = Inches(5.933)
    erd_frame_h = Inches(4.6)
    
    # Outer Frame
    erd_frame = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, erd_frame_x, erd_frame_y, erd_frame_w, erd_frame_h)
    erd_frame.fill.solid()
    erd_frame.fill.fore_color.rgb = COLOR_WHITE
    erd_frame.line.color.rgb = COLOR_GREEN
    erd_frame.line.width = Pt(1.5)
    
    # Insert ERD Image if exists
    if os.path.exists(ERD_IMAGE_PATH):
        margin_img = Inches(0.12)
        slide5.shapes.add_picture(
            ERD_IMAGE_PATH, 
            erd_frame_x + margin_img, 
            erd_frame_y + margin_img, 
            erd_frame_w - (margin_img * 2), 
            erd_frame_h - (margin_img * 2)
        )
    else:
        # Fallback Text
        fallback_tb = slide5.shapes.add_textbox(erd_frame_x, erd_frame_y + Inches(1.8), erd_frame_w, Inches(1.0))
        tf_fb = fallback_tb.text_frame
        tf_fb.paragraphs[0].alignment = PP_ALIGN.CENTER
        run_fb = tf_fb.paragraphs[0].add_run()
        run_fb.text = "[ ERD Diagram Image Missing ]\nPlease place synctrip_clean_erd_diagram.jpg in the workspace."
        apply_font_to_run(run_fb, FONT_BODY_BOLD, 12, COLOR_GREEN)
        
    # ---------------------------------------------------------
    # SLIDE 6: Demo Script (서비스 이용 시연 흐름)
    # ---------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide6)
    add_header(slide6, "서비스 이용 시연 흐름", "유저 동작 시나리오 및 핵심 플로우", "06")
    
    # 5 Step Horizontal Roadmap Layout
    roadmap_y = Inches(2.2)
    roadmap_w = Inches(2.18)
    roadmap_h = Inches(4.3)
    roadmap_xs = [Inches(0.8), Inches(3.2), Inches(5.6), Inches(8.0), Inches(10.4)]
    
    demo_steps = [
        (
            "1단계",
            "온보딩 및 성향 입력",
            "가입 후 온보딩 단계에서 MBTI, 활동량, 가고 싶은 나라(바텀 시트 멀티 체크), 흡연 여부 등을 자유롭게 선택하고 서버에 저장합니다."
        ),
        (
            "2단계",
            "대시보드 추천 탐색",
            "메인 피드에서 AI가 호환성 점수를 기반으로 정렬해준 상대방 프로필 카드를 3:4 스와이프 뷰로 직관적으로 확인하고 매칭 카드를 넘깁니다."
        ),
        (
            "3단계",
            "나비효과 상세 분석",
            "마음에 드는 카드 클릭 시 상세 페이지로 진입하여 나와 상대방의 겹치는 여행 성향이 노란색 배지로 하이라이트된 상세 분석 리포트를 확인합니다."
        ),
        (
            "4단계",
            "실시간 대화 제안",
            "대시보드를 구경하는 동안, 다른 유저가 대화를 제안하면 화면 상단에서 떨어지는 부드러운 드롭다운 애니메이션 실시간 토스트 알림을 수신합니다."
        ),
        (
            "5단계",
            "1:1 실시간 소통",
            "제안 알림 수락 시 1:1 실시간 대화방으로 진입하며, 신원인증 배지(학교/직장 이메일 완료)를 확인하고 세부 일정을 즉시 조율합니다."
        )
    ]
    
    # Roadmap Center Connector Line
    conn_line = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), roadmap_y + Inches(0.4), Inches(11.78), Inches(0.02))
    conn_line.fill.solid()
    conn_line.fill.fore_color.rgb = COLOR_LIGHT_GREEN
    conn_line.line.fill.background()
    
    for idx, x_pos in enumerate(roadmap_xs):
        # Draw Card
        tf_demo = create_card_with_text_box(slide6, x_pos, roadmap_y + Inches(0.8), roadmap_w, roadmap_h - Inches(0.8), COLOR_WHITE, COLOR_BORDER_GRAY, border_width_pt=1.0)
        
        runs_demo = [
            {
                'runs': [
                    {'text': demo_steps[idx][1], 'font': FONT_BODY_BOLD, 'size': 12.5, 'color': COLOR_DARK_TEXT}
                ],
                'space_after': 8,
                'align': PP_ALIGN.CENTER
            },
            {
                'runs': [
                    {'text': demo_steps[idx][2], 'font': FONT_BODY_LIGHT, 'size': 9.5, 'color': COLOR_MUTED_TEXT}
                ],
                'line_spacing': 1.15
            }
        ]
        add_runs_to_text_frame(tf_demo, runs_demo)
        
        # Step Circle badge on the line
        circ = slide6.shapes.add_shape(MSO_SHAPE.OVAL, x_pos + roadmap_w / 2 - Inches(0.35), roadmap_y + Inches(0.05), Inches(0.7), Inches(0.7))
        circ.fill.solid()
        circ.fill.fore_color.rgb = COLOR_GREEN
        circ.line.fill.background()
        tf_c = circ.text_frame
        tf_c.margin_left = tf_c.margin_right = tf_c.margin_top = tf_c.margin_bottom = Inches(0.01)
        tf_c.paragraphs[0].alignment = PP_ALIGN.CENTER
        r_c = tf_c.paragraphs[0].add_run()
        r_c.text = str(idx + 1)
        apply_font_to_run(r_c, FONT_TITLE, 13, COLOR_WHITE)
        
        # Bottom small badges inside cards to visualize features
        badge_y = roadmap_y + roadmap_h - Inches(0.48)
        if idx == 0:
            b_shp = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos + Inches(0.35), badge_y, Inches(1.48), Inches(0.32))
            b_shp.fill.solid()
            b_shp.fill.fore_color.rgb = COLOR_LIGHT_GREEN
            b_shp.line.fill.background()
            tf_b = b_shp.text_frame
            tf_b.margin_left = tf_b.margin_right = tf_b.margin_top = tf_b.margin_bottom = Inches(0.01)
            tf_b.paragraphs[0].alignment = PP_ALIGN.CENTER
            r_b = tf_b.paragraphs[0].add_run()
            r_b.text = "✓ 국가/성향 입력"
            apply_font_to_run(r_b, FONT_BODY_BOLD, 7.5, COLOR_GREEN)
        elif idx == 1:
            b_shp = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos + Inches(0.35), badge_y, Inches(1.48), Inches(0.32))
            b_shp.fill.solid()
            b_shp.fill.fore_color.rgb = COLOR_LIGHT_GREEN
            b_shp.line.fill.background()
            tf_b = b_shp.text_frame
            tf_b.margin_left = tf_b.margin_right = tf_b.margin_top = tf_b.margin_bottom = Inches(0.01)
            tf_b.paragraphs[0].alignment = PP_ALIGN.CENTER
            r_b = tf_b.paragraphs[0].add_run()
            r_b.text = "✓ 3:4 스와이프"
            apply_font_to_run(r_b, FONT_BODY_BOLD, 7.5, COLOR_GREEN)
        elif idx == 2:
            b_shp = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos + Inches(0.35), badge_y, Inches(1.48), Inches(0.32))
            b_shp.fill.solid()
            b_shp.fill.fore_color.rgb = COLOR_YELLOW
            b_shp.line.fill.background()
            tf_b = b_shp.text_frame
            tf_b.margin_left = tf_b.margin_right = tf_b.margin_top = tf_b.margin_bottom = Inches(0.01)
            tf_b.paragraphs[0].alignment = PP_ALIGN.CENTER
            r_b = tf_b.paragraphs[0].add_run()
            r_b.text = "✓ 노란색 하이라이트"
            apply_font_to_run(r_b, FONT_BODY_BOLD, 7.5, COLOR_DARK_TEXT)
        elif idx == 3:
            b_shp = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos + Inches(0.35), badge_y, Inches(1.48), Inches(0.32))
            b_shp.fill.solid()
            b_shp.fill.fore_color.rgb = COLOR_GREEN
            b_shp.line.fill.background()
            tf_b = b_shp.text_frame
            tf_b.margin_left = tf_b.margin_right = tf_b.margin_top = tf_b.margin_bottom = Inches(0.01)
            tf_b.paragraphs[0].alignment = PP_ALIGN.CENTER
            r_b = tf_b.paragraphs[0].add_run()
            r_b.text = "✉ 대화 요청 토스트"
            apply_font_to_run(r_b, FONT_BODY_BOLD, 7.5, COLOR_WHITE)
        elif idx == 4:
            b_shp = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos + Inches(0.35), badge_y, Inches(1.48), Inches(0.32))
            b_shp.fill.solid()
            b_shp.fill.fore_color.rgb = COLOR_LIGHT_GREEN
            b_shp.line.fill.background()
            tf_b = b_shp.text_frame
            tf_b.margin_left = tf_b.margin_right = tf_b.margin_top = tf_b.margin_bottom = Inches(0.01)
            tf_b.paragraphs[0].alignment = PP_ALIGN.CENTER
            r_b = tf_b.paragraphs[0].add_run()
            r_b.text = "💬 1:1 실시간 대화"
            apply_font_to_run(r_b, FONT_BODY_BOLD, 7.5, COLOR_DARK_TEXT)

    # ---------------------------------------------------------
    # SLIDE 7: End (감사합니다)
    # ---------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background_white(slide7)
    
    # Overlapping circle decos (echoing Slide 1 cover)
    end_shape1 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(1.0), Inches(4.5), Inches(5.5))
    end_shape1.fill.solid()
    end_shape1.fill.fore_color.rgb = COLOR_LIGHT_GREEN
    end_shape1.line.fill.background()
    
    end_shape2 = slide7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.0), Inches(2.2), Inches(3.6), Inches(3.6))
    end_shape2.fill.solid()
    end_shape2.fill.fore_color.rgb = COLOR_GREEN
    end_shape2.line.fill.background()
    
    end_shape3 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(3.6), Inches(2.8), Inches(2.8))
    end_shape3.fill.solid()
    end_shape3.fill.fore_color.rgb = COLOR_YELLOW
    end_shape3.line.fill.background()
    
    # Ending Text Box
    end_box = slide7.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(6.5), Inches(3.5))
    tf_end = end_box.text_frame
    runs_end = [
        {
            'runs': [
                {'text': "감사합니다", 'font': FONT_TITLE, 'size': 62, 'color': COLOR_GREEN, 'bold': True}
            ],
            'space_after': 18
        },
        {
            'runs': [
                {'text': "SyncTrip  |  AI 기반 여행 동행 매칭 서비스", 'font': FONT_BODY_MEDIUM, 'size': 20, 'color': COLOR_DARK_TEXT}
            ],
            'space_after': 14
        },
        {
            'runs': [
                {'text': "발표 자료가 성공적으로 생성되었습니다. Q&A를 시작합니다.", 'font': FONT_BODY_LIGHT, 'size': 13, 'color': COLOR_MUTED_TEXT}
            ]
        }
    ]
    add_runs_to_text_frame(tf_end, runs_end)
    
    # Save Presentation
    prs.save(OUTPUT_PPTX_PATH)
    print(f"PowerPoint Presentation saved successfully to: {OUTPUT_PPTX_PATH}")

# ---------------------------------------------------------
# FONT VERIFICATION FUNCTION
# ---------------------------------------------------------
def verify_presentation_fonts():
    """Verifies that all text runs inside the presentation are using the target font family."""
    prs = Presentation(OUTPUT_PPTX_PATH)
    allowed_fonts = {FONT_TITLE, FONT_BODY_BOLD, FONT_BODY_MEDIUM, FONT_BODY_LIGHT}
    all_ok = True
    
    print("\n" + "="*70)
    print("FONT VERIFICATION LOG")
    print("="*70)
    
    for s_idx, slide in enumerate(prs.slides):
        print(f"\n[Slide {s_idx + 1}]")
        for sh_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                for p_idx, p in enumerate(shape.text_frame.paragraphs):
                    for r_idx, r in enumerate(p.runs):
                        font_name = r.font.name
                        
                        # Find EA typeface from XML
                        rPr = r.font._element
                        ea_font = None
                        for child in rPr:
                            if child.tag.endswith('ea'):
                                ea_font = child.get('typeface')
                        
                        is_latin_ok = font_name in allowed_fonts
                        is_ea_ok = ea_font in allowed_fonts
                        
                        status = "OK"
                        if not is_latin_ok or not is_ea_ok:
                            status = "WARNING (Mismatch)"
                            all_ok = False
                            
                        # Encode text to ASCII for safe console printing
                        safe_text = r.text[:20].encode('ascii', errors='ignore').decode('ascii')
                        print(f"  Shape {sh_idx+1} | Text: '{safe_text}...' | Latin Font: {font_name} (ok={is_latin_ok}) | EA Font: {ea_font} (ok={is_ea_ok}) -> {status}")
                        
    print("\n" + "="*70)
    if all_ok:
        print("VERIFICATION RESULT: SUCCESS (All fonts strictly verified as Gmarket Sans TTF!)")
    else:
        print("VERIFICATION RESULT: WARNING (Some text runs missing target font family!)")
    print("="*70)

# ---------------------------------------------------------
# PPTX TO PDF CONVERSION FUNCTION (Spire.Presentation with embedded fonts)
# ---------------------------------------------------------
def convert_pptx_to_pdf():
    """Loads the PPTX file, embeds custom Gmarket Sans TTF fonts, and exports to PDF."""
    print("Converting PPTX to PDF using Spire.Presentation...")
    pres = SpirePresentation()
    pres.LoadFromFile(OUTPUT_PPTX_PATH)
    
    # Embed all three weights of Gmarket Sans TTF to ensure accurate vector rendering in PDF
    font_files = [
        "GmarketSansTTF/GmarketSansTTFBold.ttf",
        "GmarketSansTTF/GmarketSansTTFMedium.ttf",
        "GmarketSansTTF/GmarketSansTTFLight.ttf"
    ]
    
    for font_file in font_files:
        if os.path.exists(font_file):
            pres.AddEmbeddedFont(font_file)
            print(f"  Embedded font file: {font_file}")
        else:
            print(f"  [Warning] Font file not found at: {font_file}")
            
    pres.SaveToFile(OUTPUT_PDF_PATH, SpireFileFormat.PDF)
    pres.Dispose()
    print(f"PDF saved successfully to: {OUTPUT_PDF_PATH}")

# ---------------------------------------------------------
# MAIN EXECUTION
# ---------------------------------------------------------
if __name__ == "__main__":
    generate_presentation()
    verify_presentation_fonts()
    convert_pptx_to_pdf()
